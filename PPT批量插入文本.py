from pptx import Presentation
from pptx.util import Inches, Pt
import csv
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation(r'd:\Users\sunbrown\Desktop\胸牌模板 - 副本.pptx')
pages = 0

w = Inches(3.3074)
h = Inches(4.8819)
left = 0
top = 0.43*h
width = w
height = 0.33*h


# names = []
filename = r'd:\Users\sunbrown\Desktop\名单 - 副本.csv'
f = open(filename, encoding='UTF-8')
reader = csv.reader(f)
for slide, row in zip(prs.slides, reader):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.paragraphs[0].text = row[0]  # 插入第一段文本
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.color.rgb = RGBColor(32, 56, 100)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    str1 = row[1].split('#')
    for s in str1:
        p = tf.add_paragraph()  # 插入第二段文本
        p.text = s
        p.font.bold = True
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(32, 56, 100)
        p.alignment = PP_ALIGN.CENTER

    # tf.fit_text(max_size=Pt(40))  # 自动设置文本字体大小，原模块不支持中文，需要更改


prs.save(r"d:\Users\sunbrown\Desktop\result.pptx")
f.close()
